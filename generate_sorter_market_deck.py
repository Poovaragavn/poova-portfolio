"""Generate the Sorter Market pitch deck (20 slides, 16:9) as a .pptx file.

Soft / neumorphic light theme: light gray background, white rounded cards
with soft shadows, pastel accents, dark slate text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- theme ---
BG = RGBColor(0xE9, 0xED, 0xF3)          # soft light gray-blue
CARD = RGBColor(0xFF, 0xFF, 0xFF)        # white cards
CARD_SOFT = RGBColor(0xF4, 0xF7, 0xFB)   # barely-tinted card
SHADOW = RGBColor(0xD3, 0xDA, 0xE4)      # fake soft shadow
EDGE = RGBColor(0xE4, 0xE9, 0xF0)        # hairline card border

INK = RGBColor(0x2E, 0x3A, 0x48)         # dark slate headings
BODY = RGBColor(0x51, 0x5E, 0x6E)        # body text
MUTED = RGBColor(0x8A, 0x94, 0xA3)       # captions / footer

# pastel accents (like the reference dashboard)
BLUE = RGBColor(0x9C, 0xC0, 0xEE)
BLUE_DEEP = RGBColor(0x5B, 0x8D, 0xD9)
MINT = RGBColor(0xA8, 0xD9, 0xC2)
PEACH = RGBColor(0xF6, 0xC8, 0x9B)
ROSE = RGBColor(0xEF, 0xB4, 0xB4)
LAV = RGBColor(0xC6, 0xB6, 0xEC)
PASTELS = [BLUE, MINT, PEACH, ROSE, LAV]

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------- helpers ---
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def _rounded(slide, x, y, w, h, radius):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.shadow.inherit = False
    return sp


def soft_card(slide, x, y, w, h, fill=CARD, radius=0.12, line=None):
    """White rounded card with a fake soft drop-shadow behind it."""
    sh = _rounded(slide, Emu(int(x + Inches(0.045))), Emu(int(y + Inches(0.07))),
                  w, h, radius)
    sh.fill.solid()
    sh.fill.fore_color.rgb = SHADOW
    sh.line.fill.background()
    sp = _rounded(slide, x, y, w, h, radius)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.color.rgb = EDGE
        sp.line.width = Pt(0.75)
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    return sp


def bar(slide, x, y, w, h, color, radius=0.5):
    sp = _rounded(slide, x, y, w, h, radius)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp


def add_text(slide, x, y, w, h, runs, size=18, color=BODY, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            text, opts = para
        else:
            text, opts = para, {}
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = FONT
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.italic = opts.get("italic", False)
        f.color.rgb = opts.get("color", color)
    return tb


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(title, kicker=None):
    """Standard soft-theme content slide."""
    slide = add_slide()
    set_bg(slide)
    if kicker:
        add_text(slide, Inches(0.75), Inches(0.34), Inches(11.9), Inches(0.35),
                 kicker.upper(), size=12.5, color=BLUE_DEEP, bold=True)
        title_y = Inches(0.64)
    else:
        title_y = Inches(0.45)
    add_text(slide, Inches(0.75), title_y, Inches(11.9), Inches(0.8),
             title, size=29, color=INK, bold=True)
    bar(slide, Inches(0.78), Inches(1.3), Inches(1.0), Pt(4.5), BLUE)
    bar(slide, Inches(1.86), Inches(1.3), Inches(0.28), Pt(4.5), PEACH)
    add_text(slide, Inches(0.75), Inches(7.08), Inches(6), Inches(0.3),
             "Sorter Market — AI-Powered Marketplace for Industrial Sorting Machines",
             size=9, color=MUTED)
    add_text(slide, Inches(12.35), Inches(7.08), Inches(0.6), Inches(0.3),
             str(len(prs.slides)), size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    return slide


def section_header(slide, x, y, w, text, color=INK):
    bar(slide, x, y + Inches(0.06), Inches(0.09), Inches(0.3), BLUE_DEEP)
    add_text(slide, x + Inches(0.2), y, w - Inches(0.2), Inches(0.45),
             text, size=18, color=color, bold=True)


def bullet_list(slide, x, y, w, items, size=15, gap=0.52, marker="dot"):
    for i, item in enumerate(items):
        iy = y + Inches(gap * i)
        if marker == "check":
            add_text(slide, x, iy - Inches(0.05), Inches(0.35), Inches(0.4),
                     "✔", size=size, color=RGBColor(0x4F, 0xA8, 0x7F), bold=True)
        elif marker == "cross":
            add_text(slide, x, iy - Inches(0.05), Inches(0.35), Inches(0.4),
                     "✖", size=size, color=RGBColor(0xD2, 0x7B, 0x7B), bold=True)
        else:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.05),
                                         iy + Inches(0.12), Inches(0.11),
                                         Inches(0.11))
            dot.fill.solid()
            dot.fill.fore_color.rgb = PASTELS[i % len(PASTELS)]
            dot.line.fill.background()
            dot.shadow.inherit = False
        add_text(slide, x + Inches(0.38), iy, w - Inches(0.38), Inches(0.45),
                 item, size=size, color=BODY)


def card_grid(slide, items, x, y, w, h, cols, size=15, gap=0.24,
              pastel_dots=True, pastel_fill=False, fill=CARD, text_color=INK):
    rows = (len(items) + cols - 1) // cols
    gx = Inches(gap)
    gy = Inches(gap)
    cw = Emu(int((w - gx * (cols - 1)) / cols))
    ch = Emu(int((h - gy * (rows - 1)) / rows))
    for i, item in enumerate(items):
        r, c = divmod(i, cols)
        cx = Emu(int(x + c * (cw + gx)))
        cy = Emu(int(y + r * (ch + gy)))
        card_fill = PASTELS[i % len(PASTELS)] if pastel_fill else fill
        soft_card(slide, cx, cy, cw, ch, fill=card_fill)
        if pastel_dots and not pastel_fill:
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Emu(int(cx + Inches(0.22))),
                                       Emu(int(cy + ch / 2 - Inches(0.07))),
                                       Inches(0.14), Inches(0.14))
            d.fill.solid()
            d.fill.fore_color.rgb = PASTELS[i % len(PASTELS)]
            d.line.fill.background()
            d.shadow.inherit = False
            add_text(slide, Emu(int(cx + Inches(0.48))), cy,
                     Emu(int(cw - Inches(0.6))), ch, item, size=size,
                     color=text_color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        else:
            add_text(slide, cx, cy, cw, ch, item, size=size, color=text_color,
                     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def two_columns(slide, left_title, left_items, right_title, right_items,
                left_marker="dot", right_marker="dot", y=Inches(1.65),
                h=Inches(5.1), left_w=Inches(6.7), gap=Inches(0.5),
                right_tint=None):
    lx = Inches(0.75)
    rx = Emu(int(lx + left_w + gap))
    rw = Emu(int(SLIDE_W - Inches(0.75) - rx))
    soft_card(slide, lx, y, left_w, h)
    soft_card(slide, rx, y, rw, h, fill=right_tint or CARD)
    pad = Inches(0.35)
    section_header(slide, lx + pad, y + Inches(0.3), left_w - pad, left_title)
    bullet_list(slide, lx + pad, y + Inches(0.95), left_w - Inches(0.7),
                left_items, marker=left_marker)
    section_header(slide, rx + pad, y + Inches(0.3), Emu(int(rw - pad)), right_title)
    bullet_list(slide, rx + pad, y + Inches(0.95), Emu(int(rw - Inches(0.7))),
                right_items, marker=right_marker)


# ================================================================ SLIDE 1 ==
s = add_slide()
set_bg(s)
# big soft hero panel
soft_card(s, Inches(0.7), Inches(0.7), Inches(11.93), Inches(6.1),
          fill=CARD, radius=0.08)
bar(s, Inches(1.2), Inches(1.75), Inches(1.5), Pt(5), BLUE)
bar(s, Inches(2.78), Inches(1.75), Inches(0.4), Pt(5), PEACH)
add_text(s, Inches(1.2), Inches(1.05), Inches(10.5), Inches(0.45),
         "PITCH DECK  •  2026", size=13, color=BLUE_DEEP, bold=True)
add_text(s, Inches(1.2), Inches(1.95), Inches(10.9), Inches(1.2),
         "Sorter Market", size=58, color=INK, bold=True)
add_text(s, Inches(1.2), Inches(3.05), Inches(10.9), Inches(0.6),
         "AI-Powered Marketplace for Industrial Sorting Machines",
         size=22, color=MUTED)
# tagline pills
tags = [("Buy", BLUE), ("Sell", MINT), ("Rent", PEACH), ("Service", ROSE),
        ("Finance", LAV)]
tw = Inches(1.35)
tgap = Inches(0.25)
total = Emu(int(tw * 5 + tgap * 4))
tx = Emu(int(Inches(1.2)))
for i, (label, col) in enumerate(tags):
    cx = Emu(int(tx + i * (tw + tgap)))
    bar(s, cx, Inches(4.15), tw, Inches(0.55), col, radius=0.5)
    add_text(s, cx, Inches(4.15), tw, Inches(0.55), label, size=15,
             color=INK, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.2), Inches(4.95), Inches(10.5), Inches(0.45),
         "All in One Platform", size=16, color=BODY, bold=True)
add_text(s, Inches(1.2), Inches(6.95), Inches(10.5), Inches(0.35),
         "Connecting manufacturers, dealers, buyers & service providers",
         size=11, color=MUTED)
add_notes(s, "Sorter Market is a dedicated digital marketplace that connects "
             "manufacturers, dealers, buyers, and service providers for "
             "industrial sorting machines.")

# ================================================================ SLIDE 2 ==
s = content_slide("Executive Summary", kicker="What is Sorter Market?")
card_grid(s, [
    "AI-powered B2B marketplace",
    "Buy & sell new / used machines",
    "Service & spare parts support",
    "Finance & logistics integration",
    "One complete industrial ecosystem",
], Inches(0.75), Inches(1.9), Inches(11.85), Inches(4.6), cols=2, size=16)
add_notes(s, "We are building a single platform where industries can manage "
             "their complete machine lifecycle.")

# ================================================================ SLIDE 3 ==
s = content_slide("Business Problem", kicker="Current Challenges")
card_grid(s, [
    "No dedicated marketplace",
    "Offline buying process",
    "Lack of trusted sellers",
    "No price transparency",
    "Difficult service support",
    "High broker dependency",
], Inches(0.75), Inches(1.9), Inches(11.85), Inches(4.6), cols=3, size=15,
          pastel_dots=False, pastel_fill=False, text_color=BODY)
add_notes(s, "Today's machine market is fragmented. Buyers and sellers waste "
             "time, money, and effort.")

# ================================================================ SLIDE 4 ==
s = content_slide("Industry Overview")
section_header(s, Inches(0.75), Inches(1.6), Inches(6), "Industries We Serve")
card_grid(s, ["Coffee", "Rice", "Tea", "Cashew", "Pulses", "Seeds",
              "Food Processing", "Plastic Recycling", "Mining"],
          Inches(0.75), Inches(2.2), Inches(7.3), Inches(4.2), cols=3,
          size=13, pastel_dots=False)
section_header(s, Inches(8.5), Inches(1.6), Inches(4.1), "Market Trend")
bullet_list(s, Inches(8.55), Inches(2.35), Inches(4.1),
            ["Industrial Automation", "AI Adoption", "Smart Manufacturing"],
            size=15, gap=0.62)
add_notes(s, "Every industry moving toward automation requires reliable "
             "sorting equipment.")

# ================================================================ SLIDE 5 ==
s = content_slide("India Market Research")
two_columns(s,
            "Why India?",
            ["Growing manufacturing sector", "Make in India initiative",
             "Rising automation demand", "Expanding MSMEs",
             "Large food processing industry"],
            "Opportunity",
            ["Thousands of factories", "Growing demand for sorting machines"],
            right_tint=CARD_SOFT)
add_notes(s, "India is one of the fastest-growing markets for industrial automation.")

# ================================================================ SLIDE 6 ==
s = content_slide("Global Market Research")
two_columns(s,
            "Global Trends",
            ["AI-driven automation", "Industry 4.0", "Smart factories",
             "Vision-based inspection", "Predictive maintenance"],
            "Opportunity",
            ["Expanding international market", "Digital procurement growth"],
            right_tint=CARD_SOFT)
add_notes(s, "The global industry is moving toward intelligent and connected "
             "manufacturing.")

# ================================================================ SLIDE 7 ==
s = content_slide("Competitor Analysis")
two_columns(s,
            "Existing Platforms",
            ["IndiaMART", "TradeIndia", "Machinio", "Exapro", "Local Dealers"],
            "Our Advantage",
            ["AI Features", "Dedicated Marketplace", "Finance", "Service",
             "Spare Parts"],
            left_marker="cross", right_marker="check")
add_notes(s, "Our competitors sell machines. We provide a complete ecosystem.")

# ================================================================ SLIDE 8 ==
s = content_slide("Market Gap & Opportunity")
two_columns(s,
            "Market Gap",
            ["No specialized platform", "No AI support", "No verified listings",
             "No end-to-end services"],
            "Opportunity",
            ["First-mover advantage", "Huge untapped market",
             "Recurring revenue model"],
            left_marker="cross", right_marker="check", right_tint=CARD_SOFT)
add_notes(s, "We are solving problems that current marketplaces don't address.")

# ================================================================ SLIDE 9 ==
s = content_slide("Sorter Market Solution", kicker="Our Platform Offers")
card_grid(s, [
    "Buy Machines", "Sell Machines", "Rent Machines", "Spare Parts",
    "Service Booking", "Finance", "AI Recommendations",
], Inches(0.75), Inches(1.95), Inches(11.85), Inches(4.5), cols=4, size=15)
add_notes(s, "Everything required for industrial machines is available in one place.")

# =============================================================== SLIDE 10 ==
s = content_slide("Platform Ecosystem")
two_columns(s,
            "Connected Stakeholders",
            ["Manufacturers", "Dealers", "Buyers", "Service Providers",
             "Logistics", "Banks", "Admin"],
            "Platform Benefits",
            ["Faster transactions", "Better collaboration", "Improved trust"],
            right_marker="check")
add_notes(s, "We connect the complete industrial ecosystem on a single digital platform.")

# =============================================================== SLIDE 11 ==
s = content_slide("Business Model")
two_columns(s,
            "Business Model",
            ["B2B Marketplace", "SaaS Platform", "Service Marketplace",
             "Finance Marketplace", "Advertising Platform"],
            "Value",
            ["Digital Transformation", "Recurring Revenue"],
            right_tint=CARD_SOFT)
add_notes(s, "Our business combines marketplace, software, and service models.")

# =============================================================== SLIDE 12 ==
s = content_slide("Revenue Streams", kicker="Revenue Sources")
card_grid(s, [
    "Subscription Plans", "Listing Fees", "Sales Commission", "Advertisements",
    "Lead Generation", "Finance Commission", "AMC", "Spare Parts",
], Inches(0.75), Inches(1.95), Inches(11.85), Inches(4.5), cols=4, size=15)
add_notes(s, "We generate revenue from multiple channels, making the business "
             "sustainable.")

# =============================================================== SLIDE 13 ==
s = content_slide("Subscription Plans")
tiers = [
    ("BASIC", "Free", ["Limited Listings"], CARD, INK, False),
    ("PROFESSIONAL", "Paid Plan", ["Unlimited Listings", "Analytics",
     "Priority Support"], BLUE, INK, True),
    ("ENTERPRISE", "Custom", ["Custom Features", "API Integration",
     "Dedicated Manager"], CARD, INK, False),
]
tx = Inches(0.75)
twidth = Inches(11.85)
gapx = Inches(0.4)
cw = Emu(int((twidth - gapx * 2) / 3))
for i, (name, price, feats, fill, tcolor, highlight) in enumerate(tiers):
    cx = Emu(int(tx + i * (cw + gapx)))
    cy = Inches(1.9)
    ch = Inches(4.5)
    soft_card(s, cx, cy, cw, ch, fill=fill, radius=0.09)
    if highlight:
        bar(s, Emu(int(cx + cw / 2 - Inches(0.95))), cy + Inches(0.35),
            Inches(1.9), Inches(0.42), CARD, radius=0.5)
        add_text(s, Emu(int(cx + cw / 2 - Inches(0.95))), cy + Inches(0.35),
                 Inches(1.9), Inches(0.42), "MOST POPULAR", size=11,
                 color=BLUE_DEEP, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        name_y = cy + Inches(0.95)
    else:
        name_y = cy + Inches(0.55)
    add_text(s, cx, name_y, cw, Inches(0.5), name, size=19, color=tcolor,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx, name_y + Inches(0.55), cw, Inches(0.4), price, size=13,
             color=(INK if highlight else MUTED), align=PP_ALIGN.CENTER)
    bar(s, Emu(int(cx + cw / 2 - Inches(0.45))), name_y + Inches(1.1),
        Inches(0.9), Pt(2.5), (CARD if highlight else BLUE), radius=0.5)
    for j, feat in enumerate(feats):
        add_text(s, Emu(int(cx + Inches(0.3))), name_y + Inches(1.45 + 0.55 * j),
                 Emu(int(cw - Inches(0.6))), Inches(0.5), feat, size=14,
                 color=tcolor, align=PP_ALIGN.CENTER)
add_notes(s, "We provide flexible plans for businesses of all sizes.")

# =============================================================== SLIDE 14 ==
s = content_slide("Functional Modules", kicker="Platform Modules")
card_grid(s, [
    "User Management", "Machine Marketplace", "Seller Dashboard",
    "Buyer Dashboard", "Finance", "Service", "Spare Parts", "Admin Panel",
    "Analytics",
], Inches(0.75), Inches(1.95), Inches(11.85), Inches(4.5), cols=3, size=14)
add_notes(s, "Each module is designed to simplify business operations.")

# =============================================================== SLIDE 15 ==
s = content_slide("Business Workflow", kicker="End-to-End Digitized Process")
steps = ["Register", "List Machine", "AI Verification", "Buyer Search",
         "Quotation", "Payment", "Delivery", "Service & AMC"]
step_cols = [BLUE, MINT, PEACH, ROSE, LAV, BLUE, MINT, PEACH]
per_row = 4
fx = Inches(0.75)
fw = Inches(11.85)
fgap = Inches(0.25)
cw = Emu(int((fw - fgap * (per_row - 1)) / per_row))
for i, step in enumerate(steps):
    r, c = divmod(i, per_row)
    cx = Emu(int(fx + c * (cw + fgap)))
    cy = Inches(2.1 + r * 2.0)
    chv = s.shapes.add_shape(MSO_SHAPE.CHEVRON, cx, cy, cw, Inches(1.0))
    chv.fill.solid()
    chv.fill.fore_color.rgb = step_cols[i]
    chv.line.fill.background()
    chv.shadow.inherit = False
    tf = chv.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = step
    run.font.name = FONT
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = INK
    add_text(s, cx, cy + Inches(1.12), cw, Inches(0.35),
             f"Step {i + 1}", size=11, color=MUTED, align=PP_ALIGN.CENTER)
add_notes(s, "The entire buying and selling process is digitized from start to finish.")

# =============================================================== SLIDE 16 ==
s = content_slide("AI & Innovation", kicker="AI Features")
card_grid(s, [
    "Smart Search", "Price Prediction", "Fraud Detection",
    "AI Chatbot", "Machine Recommendation", "Market Analytics",
], Inches(0.75), Inches(1.95), Inches(11.85), Inches(4.5), cols=3, size=16,
          pastel_fill=True, text_color=INK)
add_notes(s, "AI improves accuracy, efficiency, and customer experience.")

# =============================================================== SLIDE 17 ==
s = content_slide("Technical Architecture", kicker="Technology Stack")
stack = [
    ("Frontend", "Next.js"),
    ("Backend", "Node.js"),
    ("Database", "PostgreSQL"),
    ("Cloud", "Supabase"),
    ("AI", "OpenAI  •  Gemini"),
]
sy = Inches(1.85)
for i, (layer, tech) in enumerate(stack):
    ry = Emu(int(sy + i * Inches(0.95)))
    bar(s, Inches(1.4), ry, Inches(3.2), Inches(0.72), PASTELS[i % len(PASTELS)],
        radius=0.22)
    add_text(s, Inches(1.4), ry, Inches(3.2), Inches(0.72), layer, size=16,
             color=INK, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    soft_card(s, Inches(4.85), ry, Inches(7.0), Inches(0.72), radius=0.22)
    add_text(s, Inches(4.85), ry, Inches(7.0), Inches(0.72), tech, size=16,
             color=INK, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
add_notes(s, "Our cloud-native architecture is secure, scalable, and AI-ready.")

# =============================================================== SLIDE 18 ==
s = content_slide("Roadmap & Implementation")
phases = [
    ("PHASE 1", "Marketplace Launch"),
    ("PHASE 2", "AI Features"),
    ("PHASE 3", "Finance & Services"),
    ("PHASE 4", "Global Expansion"),
]
px = Inches(0.75)
pw = Inches(11.85)
pgap = Inches(0.35)
cw = Emu(int((pw - pgap * 3) / 4))
bar(s, Inches(0.95), Inches(3.66), Inches(11.45), Pt(3), RGBColor(0xC7, 0xD1, 0xDE))
for i, (phase, desc) in enumerate(phases):
    cx = Emu(int(px + i * (cw + pgap)))
    col = PASTELS[i % len(PASTELS)]
    soft_card(s, cx, Inches(2.3), cw, Inches(1.1), fill=col, radius=0.14)
    add_text(s, cx, Inches(2.44), cw, Inches(0.4), phase, size=13,
             color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx, Inches(2.82), cw, Inches(0.5), desc, size=13, color=INK,
             align=PP_ALIGN.CENTER)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx + cw / 2 - Inches(0.09))),
                             Inches(3.56), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BLUE_DEEP
    dot.line.fill.background()
    dot.shadow.inherit = False
    soft_card(s, cx, Inches(4.1), cw, Inches(1.4), radius=0.14)
    note = ["Core buy/sell platform goes live", "Intelligence layer added",
            "Ecosystem services enabled", "International markets entered"][i]
    add_text(s, Emu(int(cx + Inches(0.15))), Inches(4.1),
             Emu(int(cw - Inches(0.3))), Inches(1.4), note, size=11.5,
             color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_notes(s, "Our roadmap focuses on steady growth from a national platform to "
             "a global ecosystem.")

# =============================================================== SLIDE 19 ==
s = content_slide("Financial Opportunity & ROI", kicker="Business Benefits")
benefits = [
    ("Multiple Revenue Streams", "Income from subscriptions, commissions, ads & services"),
    ("Recurring Subscription Income", "Predictable SaaS-style revenue every month"),
    ("Low Operating Cost", "Digital-first platform with lean operations"),
    ("High Scalability", "Cloud-native model grows across industries & geographies"),
    ("Strong ROI", "Compounding returns as the network effect kicks in"),
]
by = Inches(1.85)
for i, (title, desc) in enumerate(benefits):
    ry = Emu(int(by + i * Inches(1.0)))
    soft_card(s, Inches(0.95), ry, Inches(11.45), Inches(0.82), radius=0.16)
    bar(s, Inches(0.95), ry, Inches(0.13), Inches(0.82),
        PASTELS[i % len(PASTELS)], radius=0.5)
    add_text(s, Inches(1.35), Emu(int(ry + Inches(0.08))), Inches(5.2),
             Inches(0.6), title, size=15, color=INK, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(6.65), Emu(int(ry + Inches(0.08))), Inches(5.55),
             Inches(0.6), desc, size=12.5, color=MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
add_notes(s, "The platform offers long-term profitability with scalable "
             "recurring revenue.")

# =============================================================== SLIDE 20 ==
s = add_slide()
set_bg(s)
soft_card(s, Inches(0.7), Inches(0.7), Inches(11.93), Inches(6.1),
          fill=CARD, radius=0.08)
add_text(s, Inches(1.2), Inches(1.05), Inches(10.9), Inches(0.45),
         "OUR VISION", size=13, color=BLUE_DEEP, bold=True)
add_text(s, Inches(1.2), Inches(1.5), Inches(10.9), Inches(1.4),
         "Become the world's leading AI-powered marketplace for industrial "
         "sorting machines.", size=27, color=INK, bold=True, line_spacing=1.15)
add_text(s, Inches(1.2), Inches(2.95), Inches(10.9), Inches(0.4),
         "FUTURE SCOPE", size=13, color=BLUE_DEEP, bold=True)
scope = ["Global Expansion", "Mobile App", "IoT Integration",
         "Predictive Maintenance", "AI Automation"]
scw = Inches(2.02)
sgap = Inches(0.18)
for i, item in enumerate(scope):
    cx = Emu(int(Inches(1.2) + i * (scw + sgap)))
    bar(s, cx, Inches(3.42), scw, Inches(0.8), PASTELS[i % len(PASTELS)],
        radius=0.28)
    add_text(s, Emu(int(cx + Inches(0.08))), Inches(3.42),
             Emu(int(scw - Inches(0.16))), Inches(0.8), item, size=11.5,
             color=INK, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
bar(s, Inches(1.2), Inches(4.75), Inches(10.9), Pt(2), EDGE, radius=0.5)
add_text(s, Inches(1.2), Inches(5.0), Inches(10.9), Inches(1.1),
         [('"Sorter Market is transforming the industrial machine marketplace',
           {"size": 17, "italic": True, "color": MUTED}),
          ('through AI, trust, and digital innovation."',
           {"size": 17, "italic": True, "color": MUTED})],
         align=PP_ALIGN.CENTER, line_spacing=1.25)
add_text(s, Inches(1.2), Inches(6.25), Inches(10.9), Inches(0.45),
         "Thank You", size=17, color=BLUE_DEEP, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.2), Inches(6.95), Inches(10.9), Inches(0.35),
         "Sorter Market — Buy  •  Sell  •  Rent  •  Service  •  Finance",
         size=10, color=MUTED)
add_notes(s, "Our goal is to create a complete digital ecosystem that "
             "simplifies every stage of the industrial equipment lifecycle.")

# ------------------------------------------------------------------ save ---
out = "Sorter_Market_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides)} slides")
